#!/usr/bin/env python3
"""
Simple PPTX -> Markdown converter.

Usage:
  python scripts/pptx_to_md.py "path/to/1.Assessment Criteria and Classroom Regulations and Course Progress Plans .pptx"

Output:
  Creates a Markdown file with the same base name and an `images/<base>/` folder for embedded images.
"""
import argparse
import os
import re
import mimetypes
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def safe_name(s: str) -> str:
    # Replace problematic characters for file names
    return re.sub(r"[^A-Za-z0-9_.()-]+", "_", s).strip('_')


def ensure_dir(p):
    if not os.path.exists(p):
        os.makedirs(p, exist_ok=True)


def image_ext_from_image(image):
    # Try to guess extension from content_type, fallback to .png
    ct = getattr(image, 'content_type', None)
    if ct:
        ext = mimetypes.guess_extension(ct)
        if ext:
            return ext
    # fallback to filename extension if present
    fname = getattr(image, 'filename', '')
    if fname and '.' in fname:
        return os.path.splitext(fname)[1]
    return '.png'


def pptx_to_md(pptx_path, md_path=None):
    prs = Presentation(pptx_path)

    base = os.path.splitext(os.path.basename(pptx_path))[0]
    base_safe = safe_name(base)

    if md_path is None:
        # Put output in same folder as PPTX
        md_path = os.path.join(os.path.dirname(pptx_path), base_safe + '.md')

    images_dir = os.path.join(os.path.dirname(md_path), 'images', base_safe)
    ensure_dir(images_dir)

    lines = []
    lines.append(f"# {base}\n")

    for i, slide in enumerate(prs.slides, start=1):
        title = None
        if slide.shapes.title and slide.shapes.title.text:
            title = slide.shapes.title.text.strip()
        lines.append(f"## Slide {i}" + (f": {title}" if title else "") + "\n")

        img_count = 0
        # Collect text shapes
        for shp in slide.shapes:
            try:
                if shp.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    img = shp.image
                    img_count += 1
                    ext = image_ext_from_image(img)
                    img_name = f"slide-{i}-img-{img_count}{ext}"
                    img_path = os.path.join(images_dir, img_name)
                    with open(img_path, 'wb') as f:
                        f.write(img.blob)
                    rel = os.path.relpath(img_path, os.path.dirname(md_path)).replace('\\', '/')
                    lines.append(f"![Slide {i} image {img_count}]({rel})\n")
                elif shp.has_text_frame:
                    text = shp.text_frame.text.strip()
                    if text:
                        # Normalize multiple blank lines
                        text = re.sub(r"\n{2,}", "\n\n", text)
                        # Add as paragraph(s)
                        for para in text.split('\n'):
                            lines.append(para)
                        lines.append('')
            except Exception as e:
                # Be resilient to odd shapes
                lines.append(f"<!-- skipped a shape due to error: {e} -->")

        # Notes
        try:
            notes = ''
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                lines.append('\n**Notes:**')
                for ln in notes.split('\n'):
                    lines.append(ln)
                lines.append('')
        except Exception:
            pass

        lines.append('---')

    with open(md_path, 'w', encoding='utf-8') as f:
        f.write('\n'.join(lines))

    return md_path


if __name__ == '__main__':
    parser = argparse.ArgumentParser(description='Convert PPTX to Markdown (extract text and images).')
    parser.add_argument('pptx', help='Path to .pptx file')
    parser.add_argument('--out', '-o', help='Output .md file path (optional)')
    args = parser.parse_args()

    pptx = os.path.abspath(args.pptx)
    out = args.out
    if out:
        out = os.path.abspath(out)

    md = pptx_to_md(pptx, md_path=out)
    print(f"Created: {md}")
